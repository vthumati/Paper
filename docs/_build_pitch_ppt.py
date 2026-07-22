"""Build the Paper investor pitch as a native PowerPoint deck.

Mirrors docs/pitch.html (same content, numbers and "ledger" visual identity:
deep-green brand + brass/seal-gold accent on paper) as a 16:9 .pptx.

Run:  py docs/_build_pitch_ppt.py
Out:  docs/Paper_Investor_Pitch.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- palette (from pitch.html light theme) ----------------------------------
BRAND = RGBColor(0x17, 0x35, 0x2A)        # deep forest — titles, hero ground
BRAND2 = RGBColor(0x28, 0x5D, 0x47)       # lighter green — secondary bars
INK = RGBColor(0x14, 0x23, 0x1D)          # near-black ink — body
PAPER = RGBColor(0xF3, 0xF5, 0xF0)        # green-tinted paper — slide ground
CARD = RGBColor(0xFF, 0xFF, 0xFF)         # raised cards
ACCENT = RGBColor(0xA6, 0x7C, 0x1A)       # brass / seal gold
ACCENT_SOFT = RGBColor(0xF0, 0xE6, 0xCF)  # gold tint — total rows, chips
MUTED = RGBColor(0x5C, 0x6B, 0x62)        # green-grey
LINE = RGBColor(0xDD, 0xE3, 0xD8)         # hairline
LINE_STRONG = RGBColor(0xC4, 0xCE, 0xBD)
WARN = RGBColor(0xB4, 0x53, 0x1F)
CLOUD = RGBColor(0xCD, 0xDC, 0xD2)        # muted text on dark ground
NEARWHITE = RGBColor(0xEE, 0xF4, 0xEE)

SERIF = "Georgia"
SANS = "Segoe UI"
MONO = "Consolas"

EMU_IN = 914400
SW, SH = 13.333, 7.5
OUT = Path(__file__).resolve().parent / "Paper_Investor_Pitch.pptx"


# --- helpers ----------------------------------------------------------------
def _font(run, size, color, bold=False, italic=False, name=SANS):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = name
    f.color.rgb = color


def text(slide, l, t, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=None, space_after=0, wrap=True):
    """paras: list of paragraphs; each paragraph = list of run tuples
    (txt, size, color, bold, italic, name)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for run_spec in para:
            txt, size, color = run_spec[0], run_spec[1], run_spec[2]
            bold = run_spec[3] if len(run_spec) > 3 else False
            italic = run_spec[4] if len(run_spec) > 4 else False
            name = run_spec[5] if len(run_spec) > 5 else SANS
            r = p.add_run()
            r.text = txt
            _font(r, size, color, bold, italic, name)
    return tb


def rect(slide, l, t, w, h, fill, line=None, line_w=1.0, rounded=False, radius=0.06):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def new_slide(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rect(s, -0.06, -0.06, SW + 0.12, SH + 0.12, bg)
    return s


def eyebrow(slide, l, t, n, label, color=ACCENT, w=8.0):
    rect(slide, l, t + 0.10, 0.34, 0.02, color)  # tick
    text(slide, l + 0.44, t - 0.05, w, 0.35,
         [[(f"{n}  —  {label.upper()}", 11.5, color, True, False, MONO)]])


def section_title(slide, l, t, title, sub=None, tw=11.9):
    text(slide, l, t, tw, 1.0, [[(title, 27, BRAND, False, False, SERIF)]],
         line_spacing=1.0)
    if sub:
        text(slide, l, t + 0.72, tw, 0.7, [[(sub, 13.5, MUTED)]], line_spacing=1.15)


def card(slide, l, t, w, h, fill=CARD, line=LINE, line_w=1.0, rounded=True, radius=0.05):
    return rect(slide, l, t, w, h, fill, line, line_w, rounded, radius)


# ===========================================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    M = 0.75  # left margin

    # ---------------- 1. COVER ------------------------------------------------
    s = new_slide(prs, BRAND)
    # brand row
    text(s, M, 0.6, 10, 0.5,
         [[("Paper", 24, NEARWHITE, True, False, SERIF),
           ("      OS FOR CORPORATE LEGAL", 11, CLOUD, False, False, MONO)]],
         anchor=MSO_ANCHOR.MIDDLE)
    seal = rect(s, SW - M - 3.1, 0.6, 3.1, 0.44, None, ACCENT, 1.0, rounded=True, radius=0.5)
    text(s, SW - M - 3.1, 0.6, 3.1, 0.44, [[("SEED ROUND · CONFIDENTIAL", 9.5, ACCENT, False, False, MONO)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # thesis
    text(s, M, 1.95, 11.5, 1.9,
         [[("One system from ", 44, NEARWHITE, False, False, SERIF),
           ("incorporation", 44, ACCENT, False, True, SERIF)],
          [("to ", 44, NEARWHITE, False, False, SERIF),
           ("exit.", 44, ACCENT, False, True, SERIF)]],
         line_spacing=1.02)

    text(s, M, 4.15, 9.7, 1.1,
         [[("Every Indian company and fund runs its cap table, compliance and fundraising on a "
            "patchwork of spreadsheets and retainer fees. Paper replaces the patchwork — on one "
            "platform, priced for India.", 15, CLOUD)]],
         line_spacing=1.28)

    # stat strip
    stats = [("~2.0 L", "DPIIT-recognised startups, +30k/yr"),
             ("₹15.7 L Cr", "AIF commitments across 1,849 funds"),
             ("~₹9,500 Cr", "TAM — corporate-legal & fund ops"),
             ("₹145 Cr", "Target ARR by Year 5 (illustrative)")]
    sw = (SW - 2 * M - 3 * 0.2) / 4
    for i, (fig, cap) in enumerate(stats):
        x = M + i * (sw + 0.2)
        rect(s, x, 5.55, sw, 1.2, RGBColor(0x1E, 0x40, 0x33), rounded=True, radius=0.06)
        text(s, x + 0.22, 5.72, sw - 0.4, 0.5, [[(fig, 21, NEARWHITE, True, False, MONO)]])
        text(s, x + 0.22, 6.22, sw - 0.4, 0.5, [[(cap, 10.5, CLOUD)]], line_spacing=1.05)

    # ---------------- 2. PROBLEM ---------------------------------------------
    s = new_slide(prs)
    eyebrow(s, M, 0.6, "01", "The problem")
    section_title(s, M, 1.0, "Corporate legal is the founder's most expensive blind spot.",
                  "Unavoidable, highly regulated work — served by fragmented, manual tooling built for advisors, not founders.")
    pains = [
        ("The cap table lives in a spreadsheet",
         "Ownership, ESOP pools and convertibles drift out of sync the moment a round closes. Founders don't know their own fully-diluted position."),
        ("Compliance is a retainer and a prayer",
         "ROC filings, board resolutions, FEMA (FC-GPR), registers — outsourced to a CS, tracked nowhere, surfaced only when a deadline is already missed."),
        ("Fundraising means re-keying everything",
         "Term sheets, PAS-4 letters, subscription agreements and the data room are rebuilt by hand each round, in Word, at ₹lakhs a raise."),
        ("Funds run on Excel too",
         "AIFs and syndicates manage LP capital accounts, drawdowns, the waterfall, NAV and Form 64C/64D by hand — the same sprawl, one layer up."),
    ]
    cw, ch = (SW - 2 * M - 0.3) / 2, 1.75
    for i, (h, b) in enumerate(pains):
        x = M + (i % 2) * (cw + 0.3)
        y = 2.5 + (i // 2) * (ch + 0.28)
        card(s, x, y, cw, ch)
        rect(s, x + 0.001, y + 0.22, 0.06, ch - 0.44, WARN)
        text(s, x + 0.3, y + 0.24, cw - 0.55, 0.5, [[(h, 16, BRAND, True, False, SERIF)]])
        text(s, x + 0.3, y + 0.72, cw - 0.55, 1.0, [[(b, 12.5, MUTED)]], line_spacing=1.18)

    # ---------------- 3. PRODUCT ---------------------------------------------
    s = new_slide(prs)
    eyebrow(s, M, 0.6, "02", "The product")
    section_title(s, M, 1.0, "The system of record for the whole company lifecycle.",
                  "Not a point tool bolted on — the workflow itself. Every artifact reads from one "
                  "source of truth; a company on Paper feeds its own numbers to the fund that backs it.")
    stages = [("STAGE 01", "Incorporate", "Wizard-driven Pvt Ltd / LLP / OPC formation with the full incorporation set drafted."),
              ("STAGE 02", "Own", "Cap table, ESOP pools & grants, founder vesting, fully-diluted & scenario views."),
              ("STAGE 03", "Raise", "Rounds & SAFEs, data room, term-sheet scanner, PAS-4 & FC-GPR on close."),
              ("STAGE 04", "Manage & fund", "Board & ROC compliance, registers — plus SEBI AIF & SPV administration with portfolio intelligence built in.")]
    rw = (SW - 2 * M - 3 * 0.22) / 4
    for i, (n, h, b) in enumerate(stages):
        x = M + i * (rw + 0.22)
        card(s, x, 2.45, rw, 1.75)
        text(s, x + 0.24, 2.62, rw - 0.4, 0.3, [[(n, 10.5, ACCENT, False, False, MONO)]])
        text(s, x + 0.24, 2.92, rw - 0.4, 0.4, [[(h, 17, BRAND, False, False, SERIF)]])
        text(s, x + 0.24, 3.36, rw - 0.45, 0.9, [[(b, 11.5, MUTED)]], line_spacing=1.16)

    pillars = [("For founders", ["Plan-based workspace — Starter → Growth → Scale",
                                 "Cap table, ESOP & 409A-style valuations",
                                 "Diligence readiness score before you raise"]),
               ("For funds & SPVs", ["Capital accounts, waterfall, NAV & 64C/64D",
                                     "KPI collection, risk signals & tear sheets",
                                     "Deal-flow & LP CRM w/ relationship intelligence"]),
               ("For investors & LPs", ["One portal across every holding",
                                        "Statements, quarterly LP reports & call notices",
                                        "Founders report KPIs from the same portal"])]
    pw = (SW - 2 * M - 2 * 0.3) / 3
    for i, (h, items) in enumerate(pillars):
        x = M + i * (pw + 0.3)
        card(s, x, 4.5, pw, 2.35)
        text(s, x + 0.26, 4.7, pw - 0.5, 0.4, [[(h, 15, BRAND, True, False, SERIF)]])
        rect(s, x + 0.26, 5.12, 0.5, 0.02, ACCENT)
        text(s, x + 0.26, 5.26, pw - 0.5, 1.5,
             [[("•  " + it, 11.5, MUTED)] for it in items], line_spacing=1.1, space_after=5)

    # ---------------- 4. WHY NOW ---------------------------------------------
    s = new_slide(prs)
    eyebrow(s, M, 0.6, "03", "Why now")
    section_title(s, M, 1.0, "The base of formalised companies and funds just crossed the tipping point.")
    facts = [("~20k / mo", "New companies incorporated in India — a compounding base that needs a system of record from day one."),
             ("~30% CAGR", "Five-year growth in AIF commitments and registered funds — the fastest-formalising slice of private capital."),
             ("$10.5 B+", "Startup capital across ~1,500 rounds in 2025 — every one a cap-table, PAS-4 and FEMA event.")]
    fw = (SW - 2 * M - 2 * 0.3) / 3
    for i, (fig, b) in enumerate(facts):
        x = M + i * (fw + 0.3)
        card(s, x, 2.5, fw, 2.3)
        text(s, x + 0.3, 2.8, fw - 0.6, 0.7, [[(fig, 30, BRAND, True, False, MONO)]])
        rect(s, x + 0.3, 3.6, 0.6, 0.02, ACCENT)
        text(s, x + 0.3, 3.78, fw - 0.6, 1.0, [[(b, 12.5, MUTED)]], line_spacing=1.2)
    text(s, M, 5.2, 10.5, 0.8,
         [[("Digitised MCA/ROC filings, SEBI's push to shorten AIF launch timelines, and a "
            "generation of founders who expect Stripe-grade software have made the manual status "
            "quo untenable. The wedge is open now.", 13.5, MUTED)]], line_spacing=1.25)

    # ---------------- 5. MARKET ----------------------------------------------
    s = new_slide(prs)
    eyebrow(s, M, 0.6, "04", "Market size")
    section_title(s, M, 1.0, "A ~₹9,500 Cr spend pool, barely digitised.",
                  "Paper converts today's manual services spend into software + a services marketplace.")
    # bars
    bx, bw_full = M, 6.0
    bars = [("TAM", "all private cos + funds", "₹9,500 Cr", 1.0, BRAND),
            ("SAM", "digitally-addressable", "₹2,500 Cr", 0.263, BRAND2),
            ("SOM", "obtainable in 5 yrs (~6%)", "₹145 Cr", 0.06, ACCENT)]
    by = 2.7
    for name, desc, val, frac, col in bars:
        text(s, bx, by, 4.0, 0.3,
             [[(name + "   ", 13, INK, True), (desc, 10.5, MUTED, False, False, MONO)]])
        text(s, bx + bw_full - 2.0, by, 2.0, 0.3, [[(val, 13, BRAND, True, False, MONO)]],
             align=PP_ALIGN.RIGHT)
        rect(s, bx, by + 0.34, bw_full, 0.42, ACCENT_SOFT, LINE, 0.75, rounded=True, radius=0.15)
        rect(s, bx, by + 0.34, max(bw_full * frac, 0.14), 0.42, col, rounded=True, radius=0.15)
        by += 1.05
    text(s, bx, by + 0.05, bw_full, 0.6,
         [[("USD ≈ TAM $1.1B · SAM $300M · SOM $17–18M ARR. SOM is the Year-5 ARR "
            "target — a single-digit share of an under-served market.", 10.5, MUTED)]],
         line_spacing=1.15)

    # build-up table
    tx, tw = 7.15, 5.4
    rows = [("TAM build-up", "Units", "ARPU/yr", "₹ Cr"),
            ("Active private companies", "~1.6 M", "₹40k", "6,400"),
            ("New incorporations / yr", "~2.0 L", "₹12k", "240"),
            ("Cap-table & equity SaaS", "~2.3 L", "₹35k", "2,450"),
            ("AIFs & SPVs (admin)", "~7k", "₹1.6 L", "410"),
            ("Total addressable market", "", "", "~9,500")]
    _table(s, tx, 2.65, tw, rows, col_w=[2.5, 1.0, 1.0, 0.9], total_last=True)

    # ---------------- 6. BUSINESS MODEL --------------------------------------
    s = new_slide(prs)
    eyebrow(s, M, 0.6, "05", "Business model")
    section_title(s, M, 1.0, "Land on a free cap table. Expand across the lifecycle.",
                  "Product-led: founders start free, upgrade to raise, then bring their funds and advisors on. Suggested plans (INR, illustrative).")
    plans = [("STARTER", "Free", "₹0", "forever",
              ["1 entity, cap table (≤ 25 holders)", "Document vault & e-sign lite", "Getting-started guide"], False),
             ("GROWTH · MOST POPULAR", "For raising", "₹4,999/mo", "~$60/mo · billed annually",
              ["Unlimited holders + ESOP & vesting", "Rounds, SAFEs, data room, scanner", "Board resolutions + compliance calendar", "Diligence readiness & valuations"], True),
             ("SCALE", "For operating", "₹14,999/mo", "~$180/mo · billed annually",
              ["Multi-entity & scenario / waterfall", "Advanced compliance & registers", "Priority support + managed add-ons"], False)]
    pw = (SW - 2 * M - 2 * 0.28) / 3
    for i, (badge, name, price, usd, feats, feat) in enumerate(plans):
        x = M + i * (pw + 0.28)
        c = card(s, x, 2.45, pw, 3.05)
        if feat:
            c.line.color.rgb = ACCENT
            c.line.width = Pt(2.0)
        text(s, x + 0.26, 2.62, pw - 0.5, 0.3, [[(badge, 9.5, ACCENT, True, False, MONO)]])
        text(s, x + 0.26, 2.92, pw - 0.5, 0.4, [[(name, 19, BRAND, False, False, SERIF)]])
        text(s, x + 0.26, 3.42, pw - 0.5, 0.4, [[(price, 23, BRAND, True, False, MONO)]])
        text(s, x + 0.26, 3.86, pw - 0.5, 0.3, [[(usd, 10.5, MUTED, False, False, MONO)]])
        rect(s, x + 0.26, 4.24, pw - 0.52, 0.015, LINE)
        text(s, x + 0.26, 4.36, pw - 0.5, 1.1,
             [[("✓  " + f, 11, INK)] for f in feats], line_spacing=1.05, space_after=5)

    streams = [("Fund & SPV admin", "Fund from ₹1.5 L/yr · SPV ₹20k/deal"),
               ("Formation & filings", "Incorporation ₹9,999 · Managed ₹2,999/mo"),
               ("Services marketplace", "10–15% take rate on partners")]
    stw = (SW - 2 * M - 2 * 0.28) / 3
    for i, (t, v) in enumerate(streams):
        x = M + i * (stw + 0.28)
        rect(s, x, 5.75, stw, 1.05, PAPER, LINE_STRONG, 1.0, rounded=True, radius=0.06)
        text(s, x + 0.24, 5.9, stw - 0.45, 0.35, [[(t, 14, BRAND, False, False, SERIF)]])
        text(s, x + 0.24, 6.28, stw - 0.45, 0.45, [[(v, 11, ACCENT, False, False, MONO)]], line_spacing=1.05)

    # ---------------- 7. PROJECTIONS -----------------------------------------
    s = new_slide(prs)
    eyebrow(s, M, 0.6, "06", "Projections")
    section_title(s, M, 1.0, "A path to ₹145 Cr ARR in five years.",
                  "Bottoms-up from paying accounts × ARPU across four revenue lines. Illustrative estimates, not a forecast.")
    # native column chart
    cd = CategoryChartData()
    cd.categories = ["Y1", "Y2", "Y3", "Y4", "Y5"]
    cd.add_series("Total ARR (₹ Cr)", (3.0, 12.3, 35.1, 75.5, 145.0))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            Inches(M), Inches(2.55), Inches(6.1), Inches(4.15), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.gap_width = 60
    dl = plot.data_labels
    dl.number_format = '0.0'
    dl.number_format_is_linked = False
    dl.font.size = Pt(11)
    dl.font.bold = True
    dl.font.color.rgb = BRAND
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    ser = plot.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = BRAND2
    try:
        pt = ser.points[4]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = ACCENT
    except Exception:
        pass
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(11)
    ca.tick_labels.font.color.rgb = MUTED
    ca.format.line.color.rgb = LINE_STRONG
    va = chart.value_axis
    va.visible = False
    va.has_major_gridlines = False

    # projection table
    tx, tw = 7.1, 5.5
    rows = [("₹ Cr", "Y1", "Y2", "Y3", "Y4", "Y5"),
            ("Startup subscriptions", "2.1", "8.6", "25.2", "55.1", "108.0"),
            ("Fund (AIF) SaaS", "0.2", "0.8", "1.9", "3.5", "5.8"),
            ("SPV / syndicate", "0.1", "0.4", "1.0", "1.9", "3.2"),
            ("Services & marketplace", "0.6", "2.5", "7.0", "15.0", "28.0"),
            ("Total ARR", "3.0", "12.3", "35.1", "75.5", "145.0")]
    _table(s, tx, 2.55, tw, rows, col_w=[2.1, 0.68, 0.68, 0.68, 0.68, 0.68], total_last=True, fs=11)
    metrics = [("~120%", "Net revenue retention"), ("~5:1", "LTV : CAC at scale"),
               ("~18k", "Paying startups by Y5"), ("~6%", "Share of SAM")]
    mw = (tw - 0.28) / 2
    for i, (m, l) in enumerate(metrics):
        x = tx + (i % 2) * (mw + 0.28)
        y = 5.35 + (i // 2) * 0.82
        rect(s, x, y, mw, 0.72, PAPER, LINE, 0.75, rounded=True, radius=0.08)
        text(s, x + 0.18, y + 0.09, mw - 0.3, 0.35, [[(m, 16, BRAND, True, False, MONO)]])
        text(s, x + 0.18, y + 0.44, mw - 0.3, 0.25, [[(l, 10, MUTED)]])

    # ---------------- 8. POSITIONING -----------------------------------------
    s = new_slide(prs)
    eyebrow(s, M, 0.6, "07", "Why Paper wins")
    section_title(s, M, 1.0, "Global tools aren't built or priced for India. Local tools are point solutions.")
    pos = [("Carta / Pulley / Ledgy", "World-class cap tables, but US-priced and US-regulated — no PAS-4, no FC-GPR, no ROC, no SEBI AIF."),
           ("Fund-ops tools (global)", "Vestberry, 4Degrees, Carta fund admin — portfolio intelligence, deal CRM and fund accounting as three more dollar subscriptions, none SEBI-native."),
           ("Point tools (India)", "Solve one slice — a cap table, or e-sign, or a tracker. The founder still stitches five vendors together."),
           ("CS / CA retainers", "The default today: reliable but manual, opaque and expensive, with company data trapped in someone's inbox."),
           ("Paper", "One India-native platform across the entire lifecycle — company and fund — with compliance and portfolio intelligence built in.")]
    y = 2.35
    for i, (who, gap) in enumerate(pos):
        highlight = who == "Paper"
        if highlight:
            rect(s, M - 0.15, y - 0.05, SW - 2 * M + 0.3, 0.82, ACCENT_SOFT, rounded=True, radius=0.08)
        text(s, M + 0.05, y + 0.12, 3.0, 0.6, [[(who, 15, BRAND, True, False, SERIF)]])
        text(s, M + 3.3, y + 0.03, SW - 2 * M - 3.3, 0.75,
             [[(gap, 12, INK if highlight else MUTED)]], line_spacing=1.12,
             anchor=MSO_ANCHOR.MIDDLE)
        if not highlight:
            rect(s, M + 0.05, y + 0.78, SW - 2 * M - 0.1, 0.012, LINE)
        y += 0.9

    # ---------------- 9. THE ASK ---------------------------------------------
    s = new_slide(prs, BRAND)
    eyebrow(s, M, 0.7, "08", "The ask", color=ACCENT)
    text(s, M, 1.2, 11.5, 1.3,
         [[("Raising a seed round to own the system of record", 30, NEARWHITE, False, False, SERIF)],
          [("for Indian corporate legal.", 30, NEARWHITE, False, False, SERIF)]],
         line_spacing=1.05)
    text(s, M, 2.85, 10.8, 1.2,
         [[("We have a working full-stack platform spanning incorporation, cap table, fundraising, "
            "governance, compliance, fund administration and fund portfolio intelligence. Capital "
            "takes us from product to distribution — turning ~2 lakh startups and 1,849 funds into "
            "paying, expanding accounts.", 15, CLOUD)]], line_spacing=1.3)
    use = [("50%", "Go-to-market: founder & fund acquisition, marketplace partnerships"),
           ("35%", "Product & compliance-engine depth across MCA / SEBI / FEMA"),
           ("15%", "Compliance, security & SOC-grade operations")]
    uw = (SW - 2 * M - 2 * 0.3) / 3
    for i, (p, x_) in enumerate(use):
        x = M + i * (uw + 0.3)
        rect(s, x, 4.5, uw, 0.03, ACCENT)
        text(s, x, 4.65, uw, 0.6, [[(p, 26, ACCENT, True, False, MONO)]])
        text(s, x, 5.3, uw, 1.0, [[(x_, 12.5, CLOUD)]], line_spacing=1.2)
    text(s, M, 6.7, 11.5, 0.4,
         [[("Paper", 15, NEARWHITE, True, False, SERIF),
           ("   ·   OS for corporate legal   ·   Confidential", 11, CLOUD, False, False, MONO)]])

    # ---------------- 10. SOURCES --------------------------------------------
    s = new_slide(prs)
    eyebrow(s, M, 0.6, "09", "Sources & notes")
    section_title(s, M, 1.0, "Sources & disclaimer")
    text(s, M, 2.2, SW - 2 * M, 2.0,
         [[("Market facts are approximate, from public data (2025–26): DPIIT-recognised "
            "startups ~1.98 lakh (PIB, Oct 2025); AIF commitments ₹15.74 lakh crore across "
            "1,849 registered funds (SEBI, Dec 2025); new incorporations ~20k/month (MCA); 2025 "
            "startup funding ~$10.5B across ~1,518 rounds (TechCrunch / Tracxn).", 13, MUTED)]],
         line_spacing=1.35)
    rect(s, M, 4.2, SW - 2 * M, 0.015, LINE)
    text(s, M, 4.45, SW - 2 * M, 2.0,
         [[("TAM / SAM / SOM, all pricing and all financial projections are illustrative "
            "management estimates prepared for discussion — not audited figures, guidance, or a "
            "forecast of actual results.", 13, INK, True)]],
         line_spacing=1.35)
    text(s, M, 5.5, 11.5, 0.4,
         [[("Paper  ·  Prepared for investor discussion  ·  Confidential", 11, MUTED, False, False, MONO)]])

    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


def _table(slide, x, y, w, rows, col_w, total_last=False, fs=12):
    """Simple styled table: first row header, optional highlighted total row."""
    n_rows, n_cols = len(rows), len(rows[0])
    scale = w / sum(col_w)
    col_w = [c * scale for c in col_w]
    rh = 0.44
    gtbl = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                  Inches(w), Inches(rh * n_rows)).table
    gtbl.first_row = False
    gtbl.horz_banding = False
    for ci, cwid in enumerate(col_w):
        gtbl.columns[ci].width = Inches(cwid)
    for ri in range(n_rows):
        gtbl.rows[ri].height = Inches(rh)
        is_head = ri == 0
        is_total = total_last and ri == n_rows - 1
        for ci in range(n_cols):
            cell = gtbl.cell(ri, ci)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if is_total:
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_SOFT
            elif is_head:
                cell.fill.solid(); cell.fill.fore_color.rgb = CARD
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = CARD
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            r = p.add_run(); r.text = str(rows[ri][ci])
            if is_head:
                _font(r, fs - 1.5, MUTED, True, False, MONO)
            elif is_total:
                _font(r, fs, BRAND, True, False, MONO if ci else SANS)
            else:
                _font(r, fs, INK if ci else MUTED, False, False, MONO if ci else SANS)
    return gtbl


if __name__ == "__main__":
    build()
